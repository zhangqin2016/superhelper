#!/usr/bin/env python3
"""Shared Office style helper for Lily's document authoring skills.

Deterministic typography for generated .docx/.pptx: OOXML viewers resolve
latin glyphs via ascii/hAnsi and CJK glyphs via eastAsia — setting only a
latin font guarantees CJK fallback drift (Word substitutes SimSun/DengXian
per machine, which is why generated documents read as "fonts all over the
place"). These helpers set the pair everywhere it matters, plus a light
default deck theme with a contrast checker.

Import from skills:
    import lily_office_style as los
    los.style_docx(doc)                    # apply to a python-docx Document
    los.style_pptx(prs)                    # apply to a python-pptx Presentation
    los.apply_light_background(slide)      # light deck default
    los.contrast_ok(fg, bg)                # WCAG ratio guard for dark designs

Self-test:  python3 lily_office_style.py --selftest
"""

DEFAULT_LATIN_FONT = "Arial"
DEFAULT_CJK_FONT = "Microsoft YaHei"  # 微软雅黑 — present on every Windows; macOS viewers substitute PingFang automatically

LIGHT_THEME = {
    "background": "FFFFFF",
    "surface": "F7F7F5",
    "text": "1F2328",
    "muted": "6B6B66",
    "accent": "6366F1",
}


# ---------------------------------------------------------------- fonts: docx

def _docx_qn():
    from docx.oxml.ns import qn
    return qn


def _set_rfonts(rpr, latin, cjk):
    qn = _docx_qn()
    rfonts = rpr.find(qn("w:rFonts"))
    if rfonts is None:
        rfonts = rpr.makeelement(qn("w:rFonts"), {})
        rpr.insert(0, rfonts)
    rfonts.set(qn("w:ascii"), latin)
    rfonts.set(qn("w:hAnsi"), latin)
    rfonts.set(qn("w:eastAsia"), cjk)
    rfonts.set(qn("w:cs"), latin)


def style_docx(doc, latin=DEFAULT_LATIN_FONT, cjk=DEFAULT_CJK_FONT):
    """Apply the latin+CJK font pair to document defaults and every style."""
    qn = _docx_qn()
    styles_el = doc.styles.element
    doc_defaults = styles_el.find(qn("w:docDefaults"))
    if doc_defaults is not None:
        rprd = doc_defaults.find(qn("w:rPrDefault"))
        if rprd is not None:
            rpr = rprd.find(qn("w:rPr"))
            if rpr is None:
                rpr = rprd.makeelement(qn("w:rPr"), {})
                rprd.append(rpr)
            _set_rfonts(rpr, latin, cjk)
    for style in doc.styles:
        try:
            el = style.element
            rpr = el.find(qn("w:rPr"))
            if rpr is None:
                rpr = el.makeelement(qn("w:rPr"), {})
                el.append(rpr)
            _set_rfonts(rpr, latin, cjk)
        except Exception:
            continue  # a style that rejects rPr edits must not block authoring
    return doc


# ---------------------------------------------------------------- fonts: pptx

def _pptx_qn():
    from pptx.oxml.ns import qn
    return qn


def apply_ea_font(run, cjk=DEFAULT_CJK_FONT, latin=DEFAULT_LATIN_FONT):
    """Set latin (a:latin) + East-Asian (a:ea) typefaces on one pptx run."""
    qn = _pptx_qn()
    rPr = run._r.get_or_add_rPr()
    if latin:
        latin_el = rPr.find(qn("a:latin"))
        if latin_el is None:
            latin_el = rPr.makeelement(qn("a:latin"), {})
            rPr.append(latin_el)
        latin_el.set("typeface", latin)
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", cjk)


def _style_text_frame(tf, latin, cjk):
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            apply_ea_font(run, cjk=cjk, latin=latin)


def style_pptx(prs, latin=DEFAULT_LATIN_FONT, cjk=DEFAULT_CJK_FONT):
    """Apply the latin+CJK font pair to every run in the presentation."""
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    _style_text_frame(shape.text_frame, latin, cjk)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            _style_text_frame(cell.text_frame, latin, cjk)
            except Exception:
                continue
    return prs


# ---------------------------------------------------------------- deck theme

def apply_light_background(slide, hex_color=LIGHT_THEME["background"]):
    """Solid light slide background — the default for generated decks."""
    from pptx.dml.color import RGBColor
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(hex_color)


def _channel_luminance(value):
    c = value / 255.0
    return c / 12.92 if c <= 0.04045 else ((c + 0.055) / 1.055) ** 2.4


def contrast_ratio(fg_hex, bg_hex):
    """WCAG 2.x contrast ratio between two RRGGBB colors."""
    def lum(hex_color):
        r, g, b = (int(hex_color[i:i + 2], 16) for i in (0, 2, 4))
        return 0.2126 * _channel_luminance(r) + 0.7152 * _channel_luminance(g) + 0.0722 * _channel_luminance(b)
    hi, lo = max(lum(fg_hex), lum(bg_hex)), min(lum(fg_hex), lum(bg_hex))
    return (hi + 0.05) / (lo + 0.05)


def contrast_ok(fg_hex, bg_hex, minimum=4.5):
    return contrast_ratio(fg_hex, bg_hex) >= minimum


# ---------------------------------------------------------------- self-test

def _selftest():
    import tempfile
    from docx import Document
    from pptx import Presentation
    from pptx.util import Inches

    doc = Document()
    doc.add_heading("季度报告 Quarterly Report", 0)
    doc.add_paragraph("中文正文 mixed with English and 12345.")
    style_docx(doc)
    qn = _docx_qn()
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
        doc.save(tmp.name)
    from docx import Document as ReopenDoc
    reopened = ReopenDoc(tmp.name)
    rpr = reopened.styles["Normal"].element.find(qn("w:rPr"))
    rfonts = rpr.find(qn("w:rFonts"))
    assert rfonts.get(qn("w:eastAsia")) == DEFAULT_CJK_FONT, "docx eastAsia must persist after save/reopen"
    assert rfonts.get(qn("w:ascii")) == DEFAULT_LATIN_FONT, "docx latin must persist after save/reopen"

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "季度回顾 Q3 Review"
    body = slide.placeholders[1].text_frame if len(slide.placeholders) > 1 else None
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1))
    box.text_frame.text = "中文要点 mixed bullets"
    apply_light_background(slide)
    style_pptx(prs)
    pqn = _pptx_qn()
    title_run = slide.shapes.title.text_frame.paragraphs[0].runs[0]
    ea = title_run._r.get_or_add_rPr().find(pqn("a:ea"))
    assert ea is not None and ea.get("typeface") == DEFAULT_CJK_FONT, "pptx run must carry a:ea typeface"

    assert contrast_ok(LIGHT_THEME["text"], LIGHT_THEME["background"]), "theme text/bg must pass AA"
    assert not contrast_ok("1F2328", "1E2761"), "dark-on-dark must fail the guard"
    print("lily_office_style selftest ok")


if __name__ == "__main__":
    import sys
    if "--selftest" in sys.argv:
        _selftest()
    else:
        print(__doc__)
