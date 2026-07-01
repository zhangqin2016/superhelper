#!/usr/bin/env python3
"""Structured document extraction for the pre-send enrichment path.

Delegates to the bundled libraries instead of hand-rolling XML parsing in JS —
python-docx (docx), openpyxl (xlsx), python-pptx (pptx), pdfplumber (pdf).
Tables are preserved as Markdown so models without native document blocks still
receive real structure, not flattened text.

PDF/image strategy is consumer-first (no torch — stays light on ordinary
laptops): digital PDFs read their text layer with pdfplumber; only pages with
no text layer (scans) are rendered with pypdfium2 and OCR'd with RapidOCR
(onnxruntime). The OCR engine is imported lazily, so digital PDFs and Office
files never pay for it.

Usage: python extract_document.py <file_path>
Emits a single JSON object on stdout: {"ok": true, "text": "..."} or
{"ok": false, "error": "..."}.

Set LILY_PDF_ENGINE=pro-pdf (or docling) only for explicit heavyweight layout
analysis. The default path stays light and local even when pro-pdf is installed.
"""

import json
import os
import sys
import zipfile
import xml.etree.ElementTree as ET


def _rows_to_markdown(rows):
    rows = [["" if c is None else str(c) for c in row] for row in rows if row is not None]
    rows = [row for row in rows if any(cell.strip() for cell in row)]
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    rows = [row + [""] * (width - len(row)) for row in rows]
    out = ["| " + " | ".join(rows[0]) + " |", "| " + " | ".join(["---"] * width) + " |"]
    for row in rows[1:]:
        out.append("| " + " | ".join(row) + " |")
    return "\n".join(out)


W_NS = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"


def _w_attr(element, name):
    return element.attrib.get(W_NS + name, "")


def _compact(value, limit=500):
    text = " ".join(str(value or "").split())
    if len(text) <= limit:
        return text
    return text[: limit - 1] + "…"


def _element_text(element):
    paragraphs = []
    for para in element.findall(f".//{W_NS}p"):
        text = "".join(t.text or "" for t in para.iter(f"{W_NS}t")).strip()
        if text:
            paragraphs.append(text)
    if paragraphs:
        return "\n".join(paragraphs)
    return "".join(t.text or "" for t in element.iter(f"{W_NS}t")).strip()


def _comment_anchor_text(document_root):
    anchors = {}
    active = []

    def walk(element):
        tag = element.tag
        if tag == f"{W_NS}commentRangeStart":
            cid = _w_attr(element, "id")
            if cid:
                active.append(cid)
            return
        if tag == f"{W_NS}commentRangeEnd":
            cid = _w_attr(element, "id")
            if cid in active:
                active.remove(cid)
            return
        if tag == f"{W_NS}t" and element.text and active:
            for cid in active:
                anchors.setdefault(cid, []).append(element.text)
        for child in list(element):
            walk(child)

    walk(document_root)
    return {cid: _compact("".join(parts), 240) for cid, parts in anchors.items()}


def extract_docx_comments(path):
    try:
        with zipfile.ZipFile(path) as docx:
            names = set(docx.namelist())
            if "word/comments.xml" not in names:
                return []
            comments_root = ET.fromstring(docx.read("word/comments.xml"))
            anchors = {}
            if "word/document.xml" in names:
                anchors = _comment_anchor_text(ET.fromstring(docx.read("word/document.xml")))
    except Exception:
        return []

    comments = []
    for item in comments_root.findall(f".//{W_NS}comment"):
        cid = _w_attr(item, "id")
        text = _element_text(item)
        if not text:
            continue
        comments.append(
            {
                "id": cid,
                "author": _w_attr(item, "author"),
                "date": _w_attr(item, "date"),
                "anchor": anchors.get(cid, ""),
                "text": _compact(text, 1200),
            }
        )
    return comments


def _format_docx_comments(comments):
    if not comments:
        return ""
    lines = ["## Comments"]
    for index, item in enumerate(comments, start=1):
        meta = [f"Comment {item.get('id') or index}"]
        if item.get("author"):
            meta.append(f"author: {item['author']}")
        if item.get("date"):
            meta.append(f"date: {item['date']}")
        lines.append(f"- {'; '.join(meta)}")
        if item.get("anchor"):
            lines.append(f"  - Anchor: {item['anchor']}")
        lines.append(f"  - Text: {item.get('text', '')}")
    return "\n".join(lines)


def extract_docx(path):
    # python-docx walked in document order keeps headings, paragraphs and
    # tables interleaved correctly — and renders tables as real Markdown
    # tables (mammoth flattens table cells into loose paragraphs).
    from docx import Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = Document(path)
    parts = []
    for child in doc.element.body.iterchildren():
        if isinstance(child, CT_P):
            para = Paragraph(child, doc)
            text = para.text.strip()
            if not text:
                continue
            style = (para.style.name or "").lower() if para.style else ""
            if style.startswith("heading"):
                level = "".join(ch for ch in style if ch.isdigit()) or "1"
                parts.append("#" * min(int(level), 6) + " " + text)
            else:
                parts.append(text)
        elif isinstance(child, CT_Tbl):
            table = Table(child, doc)
            rows = [[cell.text for cell in row.cells] for row in table.rows]
            md = _rows_to_markdown(rows)
            if md:
                parts.append(md)
    comments = _format_docx_comments(extract_docx_comments(path))
    if comments:
        parts.append(comments)
    return "\n\n".join(parts)


def extract_xlsx(path):
    import openpyxl

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        rows = list(sheet.iter_rows(values_only=True))
        table = _rows_to_markdown(rows)
        if table:
            parts.append(f"## Sheet: {sheet.title}\n\n{table}")
    wb.close()
    return "\n\n".join(parts)


def extract_pptx(path):
    from pptx import Presentation

    prs = Presentation(path)
    parts = []
    for index, slide in enumerate(prs.slides, start=1):
        lines = []
        for shape in slide.shapes:
            if shape.has_table:
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                table = _rows_to_markdown(rows)
                if table:
                    lines.append(table)
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    lines.append(text)
        if lines:
            parts.append(f"## Slide {index}\n\n" + "\n\n".join(lines))
    return "\n\n".join(parts)


_OCR_ENGINE = None


def _get_ocr():
    # Lazy + cached: importing rapidocr pulls onnxruntime/opencv, so digital
    # PDFs and Office files that never hit a scanned page pay nothing.
    global _OCR_ENGINE
    if _OCR_ENGINE is None:
        from rapidocr_onnxruntime import RapidOCR

        _OCR_ENGINE = RapidOCR()
    return _OCR_ENGINE


def _ocr(image):
    # image may be a file path or an HxWx3 numpy array (rendered PDF page).
    result, _ = _get_ocr()(image)
    if not result:
        return ""
    # result rows are [box, text, score]; reading order is left-to-right,
    # top-to-bottom as RapidOCR returns them.
    return "\n".join(line[1] for line in result)


def _ocr_pdf_pages(path, indices):
    # Render only the text-less pages with pypdfium2 (PDFium, Apache) and OCR
    # them. One PdfDocument for the whole batch. scale=2 ≈ 144 dpi — enough for
    # OCR without rendering oversized bitmaps on an ordinary laptop.
    import numpy as np
    import pypdfium2 as pdfium

    out = {}
    pdf = pdfium.PdfDocument(path)
    try:
        for i in indices:
            pil = pdf[i].render(scale=2).to_pil().convert("RGB")
            out[i] = _ocr(np.asarray(pil))
    finally:
        pdf.close()
    return out


def _pro_pdf_enabled():
    return os.environ.get("LILY_PDF_ENGINE", "").lower() in {"pro", "pro-pdf", "docling"}


def _try_pro_pdf(path):
    # Explicit upgrade only: Docling/pro-pdf is a heavyweight layout engine and
    # may load OCR/model weights or touch the network on first use. Plain
    # pre-send attachment enrichment must stay fast and bounded, so installing
    # the runtime pack alone must not change the default PDF path.
    if not _pro_pdf_enabled():
        return None
    try:
        from docling.document_converter import DocumentConverter
    except Exception:  # noqa: BLE001 — pack not installed; this is the normal base case
        return None
    try:
        result = DocumentConverter().convert(path)
        return result.document.export_to_markdown()
    except Exception as exc:  # noqa: BLE001
        sys.stderr.write(f"pro-pdf (docling) failed, using light path: {exc}\n")
        return None


def extract_pdf(path):
    # Digital PDFs (the 90% case — exports from Word/WPS) read their text layer
    # directly: milliseconds, tens of MB. Pages with no text layer are scans,
    # so they fall through to render+OCR while preserving document page order.
    pro = _try_pro_pdf(path)
    if pro is not None:
        return pro

    import pdfplumber

    page_texts = []
    empty_pages = []
    with pdfplumber.open(path) as pdf:
        for index, page in enumerate(pdf.pages):
            chunks = []
            text = page.extract_text() or ""
            if text.strip():
                chunks.append(text.strip())
            for table in page.extract_tables() or []:
                md = _rows_to_markdown(table)
                if md:
                    chunks.append(md)
            if chunks:
                page_texts.append("\n\n".join(chunks))
            else:
                page_texts.append(None)
                empty_pages.append(index)

    if empty_pages:
        for index, text in _ocr_pdf_pages(path, empty_pages).items():
            if text.strip():
                page_texts[index] = text.strip()

    return "\n\n".join(part for part in page_texts if part)


def extract_image(path):
    # Standalone scan/photo: OCR the file directly (RapidOCR reads the path).
    return _ocr(path)


IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".webp", ".gif"}

OFFICE_EXTRACTORS = {
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".pptx": extract_pptx,
}


def main():
    if len(sys.argv) != 2:
        print(json.dumps({"ok": False, "error": "USAGE"}))
        return 1
    path = sys.argv[1]
    ext = os.path.splitext(path)[1].lower()

    if ext == ".pdf":
        extractor = extract_pdf
    elif ext in IMAGE_EXTS:
        extractor = extract_image
    elif ext in OFFICE_EXTRACTORS:
        extractor = OFFICE_EXTRACTORS[ext]
    else:
        print(json.dumps({"ok": False, "error": f"UNSUPPORTED:{ext}"}))
        return 1

    try:
        text = (extractor(path) or "").strip()
    except Exception as exc:  # noqa: BLE001 — surface the cause, never crash silently
        print(json.dumps({"ok": False, "error": f"{type(exc).__name__}: {exc}"}))
        return 1
    print(json.dumps({"ok": True, "text": text}))
    return 0


if __name__ == "__main__":
    sys.exit(main())
